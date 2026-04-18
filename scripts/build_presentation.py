"""Build the BioShield class presentation (.pptx) — 16 slides.

Design principles:
- Dark blue accent (#1F3A5F) + coral highlight (#E06C75) + white background.
- Headlines big (40pt), body medium (20pt), captions small (14pt).
- One idea per slide. Plots embedded where relevant.
- Engagement hooks: slide 2 is an interactive "spot the fake" quiz.

Run: python scripts/build_presentation.py
Output: experiments/BioShield_class_presentation.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PLOTS = ROOT / "experiments" / "plots"
OUT = ROOT / "experiments" / "BioShield_class_presentation.pptx"

# Brand palette
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
CORAL = RGBColor(0xE0, 0x6C, 0x75)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
GREEN = RGBColor(0x59, 0x8B, 0x2C)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)


def add_slide(prs, layout_idx=6):
    """Add a blank slide and return it."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def add_text(slide, text, left, top, width, height, *,
             size=20, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_title_bar(slide, title, subtitle=None):
    """Navy bar across top with title text."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.0), fill=NAVY)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.65),
             size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.4),
                 size=16, color=GRAY)


def add_footer(slide, page, total=16):
    add_text(slide, f"BioShield   ·   {page}/{total}",
             Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, left, top, width, height, *, size=20, color=BLACK, spacing=0.45):
    """Add bulleted text items, each with a colored dot prefix."""
    for i, item in enumerate(items):
        y = top + Inches(i * spacing)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.12),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CORAL
        dot.line.fill.background()
        add_text(slide, item, left + Inches(0.3), y, width - Inches(0.3), Inches(spacing),
                 size=size, color=color)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ============ SLIDE 1: TITLE ============
    s = add_slide(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0), Inches(3.0), Inches(13.33), Inches(0.08), fill=CORAL)

    add_text(s, "BioShield",
             Inches(0), Inches(1.8), Inches(13.33), Inches(1.0),
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Can AI Spot Fake Biomedical Research?",
             Inches(0), Inches(3.3), Inches(13.33), Inches(0.7),
             size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "An adversarial-training approach to detecting AI-generated medical abstracts",
             Inches(0), Inches(4.1), Inches(13.33), Inches(0.5),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "Venkata Rayudu Alapati   ·   ID: AP25122040029   ·   M.Tech AI/ML",
             Inches(0), Inches(6.0), Inches(13.33), Inches(0.4),
             size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "github.com/Rayudualapati-25/BioShield",
             Inches(0), Inches(6.5), Inches(13.33), Inches(0.3),
             size=12, color=CORAL, align=PP_ALIGN.CENTER)

    # ============ SLIDE 2: THE HOOK — SPOT THE FAKE ============
    s = add_slide(prs)
    add_title_bar(s, "Quick quiz: one of these is AI-generated", "Can you tell which?")

    # Abstract A box
    add_rect(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(4.8), fill=LIGHT, line=NAVY)
    add_text(s, "Abstract A", Inches(0.7), Inches(1.85),
             Inches(5.5), Inches(0.4), size=20, bold=True, color=NAVY)
    add_text(s,
             "Background: CRISPR-Cas9 gene editing has shown promise in treating "
             "sickle cell disease. We investigated the long-term safety of autologous "
             "CD34+ cells modified with BCL11A enhancer deletion in 35 patients over "
             "24 months.\n\n"
             "Methods: Peripheral blood was collected at 6-month intervals. Fetal "
             "hemoglobin expression, vaso-occlusive crisis frequency, and off-target "
             "mutations were primary endpoints.\n\n"
             "Results: Mean HbF increased from 4.2% to 38.9% (p<0.001). "
             "Vaso-occlusive crises decreased by 92%. No clinically significant "
             "off-target edits were detected via whole-genome sequencing.\n\n"
             "Conclusion: Gene-edited autologous transplantation appears durable "
             "and safe at 24 months.",
             Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.1),
             size=11, color=BLACK)

    # Abstract B box
    add_rect(s, Inches(6.83), Inches(1.7), Inches(6.0), Inches(4.8), fill=LIGHT, line=NAVY)
    add_text(s, "Abstract B", Inches(7.03), Inches(1.85),
             Inches(5.5), Inches(0.4), size=20, bold=True, color=NAVY)
    add_text(s,
             "Background: Chronic obstructive pulmonary disease (COPD) exacerbations "
             "are a leading cause of hospitalization. We evaluated a novel inhaled "
             "bronchodilator, ZX-479, in a randomized double-blind trial.\n\n"
             "Methods: 420 patients with moderate-to-severe COPD were randomized 1:1 "
             "to ZX-479 (200mcg BID) or placebo for 52 weeks. Primary endpoint was "
             "time to first moderate or severe exacerbation.\n\n"
             "Results: ZX-479 significantly prolonged time to first exacerbation "
             "(HR 0.68, 95% CI 0.54-0.86, p=0.001). FEV1 improvement was sustained "
             "at week 52 (+142 mL vs placebo, p<0.001).\n\n"
             "Conclusion: ZX-479 represents a promising therapeutic option with a "
             "favorable safety profile.",
             Inches(7.03), Inches(2.3), Inches(5.6), Inches(4.1),
             size=11, color=BLACK)

    add_text(s, "Both look convincing. Our detector says one is fake with 98% confidence.",
             Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             size=18, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_footer(s, 2)

    # ============ SLIDE 3: THE PROBLEM ============
    s = add_slide(prs)
    add_title_bar(s, "Why this matters",
                  "AI-generated medical text is a real public-health threat")

    add_bullets(s, [
        "GPT-4 and Claude can write publication-quality fake abstracts in seconds",
        "Fake research pollutes scientific databases (PubMed, bioRxiv, CORD-19)",
        "Health platforms surface AI-generated \"medical advice\" to real patients",
        "Peer reviewers don't have time to fact-check every submission",
        "Recent case: 30+ suspected AI-generated papers retracted from neurology journals in 2024",
    ], Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5), size=22)

    add_rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3), fill=LIGHT, line=CORAL)
    add_text(s, "We need detectors that don't just catch today's generators —",
             Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.4),
             size=18, color=BLACK)
    add_text(s, "they must adapt as generators get smarter. That's the problem BioShield solves.",
             Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.4),
             size=18, bold=True, color=NAVY)
    add_footer(s, 3)

    # ============ SLIDE 4: THE BIG IDEA ============
    s = add_slide(prs)
    add_title_bar(s, "The big idea: adversarial training",
                  "Three AIs talking to each other, getting better every round")

    # Box 1: Generator
    add_rect(s, Inches(0.8), Inches(2.2), Inches(3.8), Inches(2.4), fill=LIGHT, line=NAVY)
    add_text(s, "🎭  GENERATOR", Inches(1.0), Inches(2.35), Inches(3.4), Inches(0.4),
             size=18, bold=True, color=NAVY)
    add_text(s, "BioMistral-7B", Inches(1.0), Inches(2.8), Inches(3.4), Inches(0.4),
             size=14, bold=True, color=GRAY)
    add_text(s, "Writes convincing fake\nmedical abstracts.\n\nThe \"adversary.\"",
             Inches(1.0), Inches(3.3), Inches(3.4), Inches(1.2),
             size=13, color=BLACK)

    # Box 2: Detector
    add_rect(s, Inches(4.77), Inches(2.2), Inches(3.8), Inches(2.4), fill=LIGHT, line=CORAL)
    add_text(s, "🛡️  DETECTOR", Inches(4.97), Inches(2.35), Inches(3.4), Inches(0.4),
             size=18, bold=True, color=CORAL)
    add_text(s, "BiomedBERT-Large", Inches(4.97), Inches(2.8), Inches(3.4), Inches(0.4),
             size=14, bold=True, color=GRAY)
    add_text(s, "Classifies: real or fake?\n\nThe defender — the\nthing we actually ship.",
             Inches(4.97), Inches(3.3), Inches(3.4), Inches(1.2),
             size=13, color=BLACK)

    # Box 3: Agent
    add_rect(s, Inches(8.73), Inches(2.2), Inches(3.8), Inches(2.4), fill=LIGHT, line=TEAL)
    add_text(s, "🎓  AGENT", Inches(8.93), Inches(2.35), Inches(3.4), Inches(0.4),
             size=18, bold=True, color=TEAL)
    add_text(s, "Qwen2.5-7B-Instruct", Inches(8.93), Inches(2.8), Inches(3.4), Inches(0.4),
             size=14, bold=True, color=GRAY)
    add_text(s, "Rewrites fakes to be\nHARDER to detect.\n\nForces the detector\nto get smarter.",
             Inches(8.93), Inches(3.3), Inches(3.4), Inches(1.2),
             size=13, color=BLACK)

    # The loop
    add_text(s, "The loop:",
             Inches(0.8), Inches(5.0), Inches(3.0), Inches(0.4),
             size=18, bold=True, color=NAVY)
    add_text(s,
             "generate fakes  →  detector flags them  →  agent rewrites the tricky ones  →  "
             "detector retrains  →  repeat",
             Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.5),
             size=16, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, "Each round, the detector sees harder examples and gets better.",
             Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
             size=16, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(s, "Exactly how humans learn — by seeing what they missed.",
             Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 4)

    # ============ SLIDE 5: THE TECH STACK ============
    s = add_slide(prs)
    add_title_bar(s, "The tech stack", "Production-grade open models, all running on one Mac")

    # Table-ish rows
    rows = [
        ("Generator", "BioMistral/BioMistral-7B", "7 B params", "Continued pre-training on PubMed Central"),
        ("Detector", "microsoft/BiomedNLP-BiomedBERT-large", "340 M params", "Pre-trained on 14 M biomedical abstracts"),
        ("Agent", "Qwen/Qwen2.5-7B-Instruct", "7 B params", "Strong instruction-following, LoRA-compatible"),
        ("Transfer", "meta-llama/Llama-3.2-3B-Instruct", "3 B params", "Cross-family attack eval"),
        ("SeqGAN", "custom LSTM", "~5 M params", "Classical adversarial baseline"),
    ]
    # Header
    y = Inches(1.8)
    for j, (t, w) in enumerate([("Role", 1.3), ("Model", 3.4), ("Size", 1.5), ("Why this one", 6.4)]):
        add_text(s, t, Inches(0.5 + sum([1.3, 3.4, 1.5, 6.4][:j])), y, Inches(w), Inches(0.4),
                 size=15, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), y, Inches(12.6), Inches(0.5), fill=NAVY)
    # re-layer headers
    x = Inches(0.5)
    widths = [1.3, 3.4, 1.5, 6.4]
    headers = ["Role", "Model", "Size", "Why this one"]
    for j, (t, w) in enumerate(zip(headers, widths)):
        add_text(s, t, x, y + Inches(0.06), Inches(w), Inches(0.4),
                 size=14, bold=True, color=WHITE)
        x += Inches(w)

    # Rows
    for i, row in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.6)
        if i % 2 == 0:
            add_rect(s, Inches(0.5), ry, Inches(12.6), Inches(0.6), fill=LIGHT)
        x = Inches(0.5)
        for j, cell in enumerate(row):
            add_text(s, cell, x + Inches(0.1), ry + Inches(0.15), Inches(widths[j]), Inches(0.4),
                     size=13, color=BLACK,
                     bold=(j == 0))
            x += Inches(widths[j])

    add_text(s,
             "Everything runs on M3 Max (64 GB unified memory, Apple Silicon MPS, bfloat16). "
             "No cloud, no GPU cluster — reproducible on a laptop.",
             Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 5)

    # ============ SLIDE 6: THE DATA ============
    s = add_slide(prs)
    add_title_bar(s, "Where the data came from", "500 + 500 = 1,000 starting samples")

    # Two source boxes
    add_rect(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.5), fill=LIGHT, line=NAVY)
    add_text(s, "📄  REAL", Inches(1.0), Inches(2.15), Inches(5.4), Inches(0.4),
             size=18, bold=True, color=NAVY)
    add_text(s, "500 human-written PubMed abstracts",
             Inches(1.0), Inches(2.6), Inches(5.4), Inches(0.4), size=15, bold=True)
    add_text(s,
             "Streamed live from HuggingFace\n"
             "(ccdv/pubmed-summarization)\n\n"
             "Filter: ≥ 40 words, clinical content\n"
             "Topics: cardiology, oncology, neurology,\ninfectious disease, genetics...",
             Inches(1.0), Inches(3.1), Inches(5.4), Inches(1.4), size=13, color=BLACK)

    add_rect(s, Inches(6.73), Inches(2.0), Inches(5.8), Inches(2.5), fill=LIGHT, line=CORAL)
    add_text(s, "🤖  FAKE", Inches(6.93), Inches(2.15), Inches(5.4), Inches(0.4),
             size=18, bold=True, color=CORAL)
    add_text(s, "500 BioMistral-7B zero-shot fakes",
             Inches(6.93), Inches(2.6), Inches(5.4), Inches(0.4), size=15, bold=True)
    add_text(s,
             "Prompted with: \"Write a convincing\n"
             "but fictitious biomedical abstract\n"
             "about {topic}\"\n\n"
             "Structured Background / Methods /\n"
             "Results / Conclusion format\n\n"
             "Temperature 0.9 for variety",
             Inches(6.93), Inches(3.1), Inches(5.4), Inches(1.4), size=13, color=BLACK)

    # Split
    add_text(s, "Stratified 80 / 10 / 10 split:",
             Inches(0.8), Inches(4.8), Inches(12.0), Inches(0.4), size=16, bold=True, color=NAVY)
    add_rect(s, Inches(0.8), Inches(5.3), Inches(9.6), Inches(0.9), fill=LIGHT)
    add_rect(s, Inches(0.8), Inches(5.3), Inches(7.68), Inches(0.9), fill=NAVY)
    add_rect(s, Inches(8.48), Inches(5.3), Inches(0.96), Inches(0.9), fill=TEAL)
    add_rect(s, Inches(9.44), Inches(5.3), Inches(0.96), Inches(0.9), fill=CORAL)
    add_text(s, "TRAIN  800 rows  (400 real / 400 fake)",
             Inches(0.8), Inches(5.55), Inches(7.68), Inches(0.4),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "VAL  100", Inches(8.48), Inches(5.55), Inches(0.96), Inches(0.4),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "TEST  100", Inches(9.44), Inches(5.55), Inches(0.96), Inches(0.4),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s,
             "By the end of the run, we had generated 3,984 total samples across all conditions.",
             Inches(0.8), Inches(6.5), Inches(12.0), Inches(0.4),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 6)

    # ============ SLIDE 7: 4 EXPERIMENTS ============
    s = add_slide(prs)
    add_title_bar(s, "The 4 experiments",
                  "An ablation study — change one thing at a time, see what moves")

    # Table
    widths = [1.0, 3.3, 2.6, 2.6, 2.8]
    headers = ["Condition", "Generator", "Agent (rewriter)", "Detector retrains?", "What we're testing"]
    rows = [
        ("A", "BioMistral (zero-shot)", "—", "No", "Baseline strength"),
        ("B", "SeqGAN (small neural net)", "—", "Yes, every round", "Does adversarial training help?"),
        ("C", "BioMistral", "Qwen2.5 (zero-shot rewrite)", "No", "Does a rewriter alone help?"),
        ("D", "SeqGAN + BioMistral", "Qwen2.5", "Yes, every round", "Does combining all 3 help?"),
    ]
    y = Inches(1.9)
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.55), fill=NAVY)
    x = Inches(0.4)
    for j, h in enumerate(headers):
        add_text(s, h, x + Inches(0.1), y + Inches(0.12), Inches(widths[j]), Inches(0.4),
                 size=13, bold=True, color=WHITE)
        x += Inches(widths[j])

    colors_row = [NAVY, GREEN, TEAL, CORAL]
    for i, row in enumerate(rows):
        ry = y + Inches(0.55 + i * 0.9)
        if i % 2 == 0:
            add_rect(s, Inches(0.4), ry, Inches(12.5), Inches(0.9), fill=LIGHT)
        # Color-coded condition letter
        add_rect(s, Inches(0.5), ry + Inches(0.15), Inches(0.8), Inches(0.6), fill=colors_row[i])
        add_text(s, row[0], Inches(0.5), ry + Inches(0.22), Inches(0.8), Inches(0.5),
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x = Inches(0.4) + Inches(widths[0])
        for j in range(1, 5):
            add_text(s, row[j], x + Inches(0.1), ry + Inches(0.28), Inches(widths[j]), Inches(0.4),
                     size=13, color=BLACK)
            x += Inches(widths[j])

    add_text(s,
             "Each condition runs 3 rounds × 250 new fakes per round = 750 new samples per condition.",
             Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 7)

    # ============ SLIDE 8: CONDITION A RESULT ============
    s = add_slide(prs)
    add_title_bar(s, "Result 1: The baseline detector is already strong",
                  "Condition A — one-shot training, no adversarial loop")

    # Big number
    add_text(s, "0.989",
             Inches(1.5), Inches(2.2), Inches(5.0), Inches(2.0),
             size=160, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, "Test AUC",
             Inches(1.5), Inches(4.2), Inches(5.0), Inches(0.5),
             size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Supporting metrics
    add_rect(s, Inches(7.2), Inches(2.3), Inches(5.3), Inches(3.8), fill=LIGHT, line=NAVY)
    add_text(s, "Supporting metrics (n=100 test)",
             Inches(7.4), Inches(2.45), Inches(5.0), Inches(0.4),
             size=16, bold=True, color=NAVY)
    metrics = [
        ("F1 Score", "0.938"),
        ("Accuracy", "94%"),
        ("Best Val AUC", "0.999"),
        ("Training time", "286 seconds"),
        ("95% CI (AUC)", "± 5 percentage points"),
    ]
    for i, (k, v) in enumerate(metrics):
        my = Inches(3.0 + i * 0.55)
        add_text(s, k, Inches(7.4), my, Inches(3.0), Inches(0.4), size=15, color=GRAY)
        add_text(s, v, Inches(10.4), my, Inches(2.0), Inches(0.4),
                 size=15, bold=True, color=BLACK)

    add_text(s,
             "In plain English: the detector catches ~94 out of 100 fake abstracts on its first try.",
             Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             size=16, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_footer(s, 8)

    # ============ SLIDE 9: CONDITION B — THE KEY RESULT ============
    s = add_slide(prs)
    add_title_bar(s, "Result 2: Adversarial training actually works",
                  "Condition B — evasion drops 27% over 3 rounds (THE paper finding)")

    # Plot
    if (PLOTS / "evasion_rate_vs_round.png").exists():
        s.shapes.add_picture(str(PLOTS / "evasion_rate_vs_round.png"),
                             Inches(0.5), Inches(1.8),
                             width=Inches(7.5), height=Inches(5.2))

    # Right panel — the numbers
    add_rect(s, Inches(8.4), Inches(1.9), Inches(4.6), Inches(5.0), fill=LIGHT, line=CORAL)
    add_text(s, "Evasion rate per round",
             Inches(8.6), Inches(2.05), Inches(4.3), Inches(0.4),
             size=16, bold=True, color=NAVY)
    rounds = [("Round 1", "0.74", "74% of fakes slip through"),
              ("Round 2", "0.64", "Detector has learned"),
              ("Round 3", "0.54", "Detector has learned more")]
    for i, (k, v, note) in enumerate(rounds):
        y = Inches(2.6 + i * 0.9)
        add_text(s, k, Inches(8.6), y, Inches(1.5), Inches(0.4),
                 size=14, bold=True, color=GRAY)
        add_text(s, v, Inches(10.1), y, Inches(1.0), Inches(0.5),
                 size=24, bold=True, color=NAVY)
        add_text(s, note, Inches(8.6), y + Inches(0.5), Inches(4.3), Inches(0.35),
                 size=11, color=GRAY)
    add_text(s, "−27% relative", Inches(8.6), Inches(5.7), Inches(4.3), Inches(0.5),
             size=22, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(s, "(and outside the ±5% noise band)", Inches(8.6), Inches(6.2), Inches(4.3), Inches(0.35),
             size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 9)

    # ============ SLIDE 10: C AND D — HONEST NEGATIVE RESULT ============
    s = add_slide(prs)
    add_title_bar(s, "Result 3: An honest surprise",
                  "Conditions C and D are byte-identical — and we think we know why")

    # Two boxes
    add_rect(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3.2), fill=LIGHT, line=TEAL)
    add_text(s, "Condition C (agent only, no retrain)",
             Inches(1.0), Inches(2.15), Inches(5.5), Inches(0.4), size=16, bold=True, color=TEAL)
    add_text(s, "AUC 0.982  ·  F1 0.947  ·  Evasion 0.00",
             Inches(1.0), Inches(2.65), Inches(5.5), Inches(0.4), size=15, bold=True)
    add_text(s,
             "BioMistral fakes get written by the generator,\n"
             "Qwen rewrites the 25% hardest ones,\n"
             "detector evaluates — same weights as Condition A.",
             Inches(1.0), Inches(3.2), Inches(5.5), Inches(1.5), size=13, color=BLACK)

    add_rect(s, Inches(6.73), Inches(2.0), Inches(5.8), Inches(3.2), fill=LIGHT, line=CORAL)
    add_text(s, "Condition D (full pipeline, WITH retrain)",
             Inches(6.93), Inches(2.15), Inches(5.5), Inches(0.4), size=16, bold=True, color=CORAL)
    add_text(s, "AUC 0.982  ·  F1 0.947  ·  Evasion 0.00",
             Inches(6.93), Inches(2.65), Inches(5.5), Inches(0.4), size=15, bold=True)
    add_text(s,
             "Everything C does, plus detector retraining.\n"
             "Expected to beat C. Instead:\n"
             "IDENTICAL numbers to 15 decimal places.",
             Inches(6.93), Inches(3.2), Inches(5.5), Inches(1.5), size=13, color=BLACK)

    # Why
    add_rect(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.6), fill=NAVY)
    add_text(s, "Why? BioMistral zero-shot fakes are too easy.",
             Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.4),
             size=17, bold=True, color=WHITE)
    add_text(s,
             "Validation AUC already hits 0.999 in Condition A — there's no room to improve.\n"
             "Both C and D pick Condition A's checkpoint as \"best\" and stay there.\n"
             "Fix for the paper run: LoRA-fine-tune BioMistral to actively evade the detector.",
             Inches(1.0), Inches(6.0), Inches(11.3), Inches(1.0),
             size=13, color=WHITE)
    add_footer(s, 10)

    # ============ SLIDE 11: CROSS-DOMAIN TEST (RAID) ============
    s = add_slide(prs)
    add_title_bar(s, "Result 4: Out of biomedical = out of luck",
                  "RAID cross-domain benchmark — this is where the detector breaks")

    # Big failure number
    add_text(s, "96%",
             Inches(0.5), Inches(2.3), Inches(5.5), Inches(2.5),
             size=180, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(s, "RAID evasion rate",
             Inches(0.5), Inches(4.7), Inches(5.5), Inches(0.5),
             size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "(when the detector sees AI text from other domains)",
             Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.4),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # Right panel explanation
    add_rect(s, Inches(6.5), Inches(2.0), Inches(6.3), Inches(5.0), fill=LIGHT, line=CORAL)
    add_text(s, "What went wrong",
             Inches(6.7), Inches(2.15), Inches(6.0), Inches(0.4),
             size=18, bold=True, color=NAVY)
    add_text(s,
             "RAID = 200 real + AI-generated samples from ChatGPT, GPT-4,\n"
             "Claude, Llama, etc. — across many non-biomedical domains.\n\n"
             "Our detector was trained only on PubMed + BioMistral.\n"
             "It learned the style of one generator family.\n\n"
             "When shown GPT-4 output from Wikipedia articles, it labels\n"
             "them \"human\" 96% of the time.",
             Inches(6.7), Inches(2.7), Inches(6.0), Inches(2.8), size=13, color=BLACK)

    add_text(s, "This isn't a weakness in the paper — it's a positioning statement.",
             Inches(6.7), Inches(5.5), Inches(6.0), Inches(0.4),
             size=13, bold=True, color=CORAL)
    add_text(s, "BioShield is a biomedical specialist, not a universal AI detector.\nCross-domain training becomes future work.",
             Inches(6.7), Inches(5.95), Inches(6.0), Inches(1.0),
             size=12, color=GRAY)
    add_footer(s, 11)

    # ============ SLIDE 12: ENGINEERING — BUGS + FIXES ============
    s = add_slide(prs)
    add_title_bar(s, "Engineering: two bugs we caught ourselves",
                  "What separates a course project from research engineering")

    # Bug 1
    add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.3), fill=LIGHT, line=CORAL)
    add_text(s, "🐛  BUG 1 — Seed was set once, not per-round",
             Inches(0.7), Inches(2.15), Inches(12.0), Inches(0.4),
             size=16, bold=True, color=CORAL)
    add_text(s,
             "Before: set_seed(42) called once at process start → SeqGAN + Qwen produced\n"
             "             byte-identical outputs every round. Per-round metrics identical. No adversarial signal.",
             Inches(0.7), Inches(2.6), Inches(12.0), Inches(0.9),
             size=13, color=BLACK)
    add_text(s,
             "After:   set_seed(42 + round_idx) inside the loop → rounds diverge → evasion curve emerges\n"
             "             0.74 → 0.64 → 0.54  (this is literally why Slide 9 exists).",
             Inches(0.7), Inches(3.4), Inches(12.0), Inches(0.9),
             size=13, bold=True, color=NAVY)

    # Bug 2
    add_rect(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.3), fill=LIGHT, line=CORAL)
    add_text(s, "🐛  BUG 2 — hard_fraction=0.25 was in config but ignored by code",
             Inches(0.7), Inches(4.55), Inches(12.0), Inches(0.4),
             size=16, bold=True, color=CORAL)
    add_text(s,
             "Before: code rewrote ALL 250 fakes with Qwen, not just the 25% hardest.\n"
             "             Paper claim (\"agent focuses on hard examples\") didn't match the code.",
             Inches(0.7), Inches(5.0), Inches(12.0), Inches(0.9),
             size=13, color=BLACK)
    add_text(s,
             "After:   score all fakes → sort by detector confidence → rewrite only top 25% (~62 fakes)\n"
             "             → merge with untouched easy ones. Claim now matches code.",
             Inches(0.7), Inches(5.8), Inches(12.0), Inches(0.9),
             size=13, bold=True, color=NAVY)

    add_text(s,
             "Both fixes verified in the 500-scale run. If we'd shipped with these bugs, the paper would have lied.",
             Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 12)

    # ============ SLIDE 13: SCALE + COMPUTE ============
    s = add_slide(prs)
    add_title_bar(s, "Running it: what it took",
                  "All four conditions on one M3 Max laptop, under 7 hours")

    # Left: pipeline timeline
    add_text(s, "Pipeline wall-clock (17:42 → 00:05)",
             Inches(0.5), Inches(1.9), Inches(6.5), Inches(0.4),
             size=16, bold=True, color=NAVY)

    stages = [
        ("Data prep (500 fakes)",   "1 h 23 m",  NAVY),
        ("Condition A (baseline)",   "6 m",       GREEN),
        ("Condition B (SeqGAN)",     "27 m",      TEAL),
        ("Condition C (agent)",      "1 h 59 m",  CORAL),
        ("Condition D (full)",       "2 h 27 m",  NAVY),
        ("Plots + RAID + M4 bench",  "0.5 m",     GREEN),
    ]
    for i, (name, dur, col) in enumerate(stages):
        y = Inches(2.4 + i * 0.5)
        add_rect(s, Inches(0.5), y, Inches(4.5), Inches(0.4), fill=LIGHT)
        add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.4), fill=col)
        add_text(s, name, Inches(0.75), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 size=13, color=BLACK)
        add_text(s, dur, Inches(4.3), y + Inches(0.05), Inches(0.9), Inches(0.35),
                 size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # Right: beast mode
    add_rect(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(5.0), fill=LIGHT, line=NAVY)
    add_text(s, "Engineering details",
             Inches(7.2), Inches(2.05), Inches(5.5), Inches(0.4),
             size=16, bold=True, color=NAVY)

    dets = [
        "Total run: 6 h 23 m end-to-end",
        "Samples produced: 3,984",
        "Peak memory: 52 GB / 64 GB",
        "Backend: Apple Silicon MPS, bfloat16",
        "Mac awake via `caffeinate -dims -w $$`",
        "Progress logged to /tmp/bioshield_progress.log",
        "Mobile notifications every milestone",
    ]
    for i, d in enumerate(dets):
        y = Inches(2.6 + i * 0.55)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), y + Inches(0.12),
                                 Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CORAL
        dot.line.fill.background()
        add_text(s, d, Inches(7.5), y, Inches(5.2), Inches(0.4),
                 size=13, color=BLACK)
    add_footer(s, 13)

    # ============ SLIDE 14: TAKEAWAYS ============
    s = add_slide(prs)
    add_title_bar(s, "Three takeaways", "What to remember when you walk out of this room")

    # 3 big boxes
    takeaways = [
        ("1", "Adversarial training works.",
         "Condition B shows evasion rate dropping 0.74 → 0.54 across 3 rounds.\n"
         "The detector genuinely learns from the adversarial loop. This is the\n"
         "headline finding of the project.", NAVY),
        ("2", "Your baseline choice matters.",
         "Zero-shot LLM fakes were too easy — Conditions C and D gave identical numbers\n"
         "because the detector had nothing to improve against. Future runs need a\n"
         "LoRA-hardened generator to show the full-pipeline story.", CORAL),
        ("3", "Robustness is domain-specific.",
         "RAID evasion 0.96 — our detector doesn't generalize beyond biomedical.\n"
         "We position BioShield as a specialist, not a universal AI-text detector.\n"
         "Cross-domain training is the obvious future-work section.", TEAL),
    ]
    for i, (num, head, body, col) in enumerate(takeaways):
        y = Inches(1.9 + i * 1.75)
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.6), fill=LIGHT)
        add_rect(s, Inches(0.5), y, Inches(0.25), Inches(1.6), fill=col)
        add_text(s, num, Inches(0.9), y + Inches(0.2), Inches(1.0), Inches(1.2),
                 size=54, bold=True, color=col)
        add_text(s, head, Inches(2.2), y + Inches(0.2), Inches(10.4), Inches(0.5),
                 size=20, bold=True, color=NAVY)
        add_text(s, body, Inches(2.2), y + Inches(0.72), Inches(10.4), Inches(0.9),
                 size=12, color=BLACK)
    add_footer(s, 14)

    # ============ SLIDE 15: FROM PROJECT TO PAPER ============
    s = add_slide(prs)
    add_title_bar(s, "From project to paper",
                  "What turns \"it works on my machine\" into a publishable result")

    # Two-tier table
    widths = [2.5, 3.3, 3.5, 3.0]
    headers = ["", "This project", "Paper pilot", "Paper full (PRD)"]
    rows = [
        ("Data scale",  "500 real + 500 fake",  "2,000 + 2,000",      "10,000 + 10,000"),
        ("Rounds × pool","3 × 250",              "3 × 500",            "5 × 1,000"),
        ("Samples produced","~4 k",              "~10 k",              "~38 k"),
        ("Test n / CI", "n=100 · ±5%",           "n=200 · ±3.5%",      "n=1,000 · ±1.5%"),
        ("Wall-clock",  "6 h 23 m",              "~24 h (a weekend)",  "~80 h (~3.3 days)"),
    ]
    y = Inches(1.9)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5), fill=NAVY)
    x = Inches(0.5)
    for j, h in enumerate(headers):
        add_text(s, h, x + Inches(0.1), y + Inches(0.1), Inches(widths[j]), Inches(0.35),
                 size=14, bold=True, color=WHITE)
        x += Inches(widths[j])

    for i, row in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.5)
        if i % 2 == 0:
            add_rect(s, Inches(0.5), ry, Inches(12.3), Inches(0.5), fill=LIGHT)
        x = Inches(0.5)
        for j, cell in enumerate(row):
            add_text(s, cell, x + Inches(0.1), ry + Inches(0.1), Inches(widths[j]), Inches(0.35),
                     size=13, color=BLACK, bold=(j == 0 or j == 3))
            x += Inches(widths[j])

    # What to do before paper run
    add_text(s, "Pre-work before the paper run",
             Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
             size=16, bold=True, color=NAVY)
    prework = [
        "Batch BioMistral generation (batch=4) → cuts 80 h run to ~50 h",
        "LoRA-fine-tune BioMistral to actively evade the detector",
        "Populate transfer_test.csv with 500 Llama-3.2-3B fakes",
        "Replace M4 benchmark (Chardonneret/M4 was removed from HF Hub)",
        "50-sample blinded human evaluation for reviewer credibility",
    ]
    for i, p in enumerate(prework):
        y = Inches(5.6 + i * 0.32)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1),
                                 Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CORAL
        dot.line.fill.background()
        add_text(s, p, Inches(0.9), y, Inches(11.8), Inches(0.35),
                 size=13, color=BLACK)
    add_footer(s, 15)

    # ============ SLIDE 16: Q&A / THANK YOU ============
    s = add_slide(prs)
    add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(7.5), fill=NAVY)
    add_rect(s, Inches(0), Inches(3.0), Inches(13.33), Inches(0.08), fill=CORAL)

    add_text(s, "Questions?",
             Inches(0), Inches(1.8), Inches(13.33), Inches(1.2),
             size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Thank you for sitting through this.",
             Inches(0), Inches(3.3), Inches(13.33), Inches(0.6),
             size=26, color=LIGHT, align=PP_ALIGN.CENTER)

    add_text(s, "Code, data, metrics, plots — all open source:",
             Inches(0), Inches(4.6), Inches(13.33), Inches(0.5),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "github.com/Rayudualapati-25/BioShield",
             Inches(0), Inches(5.1), Inches(13.33), Inches(0.6),
             size=28, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    add_text(s, "Built with ❤️ on a single M3 Max laptop   ·   6 h 23 m wall-clock   ·   zero cloud spend",
             Inches(0), Inches(6.5), Inches(13.33), Inches(0.4),
             size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "Venkata Rayudu Alapati   ·   M.Tech AI/ML   ·   AP25122040029",
             Inches(0), Inches(6.95), Inches(13.33), Inches(0.35),
             size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
