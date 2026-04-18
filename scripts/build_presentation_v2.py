"""Claude-design class presentation for BioShield.

Editorial aesthetic: cream backgrounds, serif hero numerals, hairline rules,
generous whitespace. Built to feel like an essay you're walking through, not
a corporate deck. 16 slides.

Run:
    python3 scripts/build_presentation_v2.py

Output:
    experiments/BioShield_class_presentation_v2.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ─── palette ──────────────────────────────────────────────────────────────────
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)  # page background
PAPER      = RGBColor(0xFA, 0xF7, 0xF2)  # slightly brighter alternate bg
INK        = RGBColor(0x1A, 0x1A, 0x1A)  # primary text, almost black
SOFT_INK   = RGBColor(0x3A, 0x3A, 0x3A)  # body text
MUTED      = RGBColor(0x8C, 0x83, 0x78)  # metadata, captions
HAIRLINE   = RGBColor(0xC8, 0xBE, 0xAF)  # thin rules
CARD_EDGE  = RGBColor(0xD8, 0xCE, 0xBE)  # card borders
CORAL      = RGBColor(0xD9, 0x77, 0x57)  # signature accent
CORAL_SOFT = RGBColor(0xEF, 0xC6, 0xB3)  # tinted fills
DEEP       = RGBColor(0x12, 0x1F, 0x30)  # dramatic dark slides
CREAM_ON_D = RGBColor(0xEE, 0xE6, 0xD8)  # cream tone for dark-bg text
FOREST     = RGBColor(0x2C, 0x5E, 0x4A)  # small accents for "learning" curve

# ─── typography ───────────────────────────────────────────────────────────────
SERIF = "Georgia"       # universally available, editorial feel
SANS  = "Helvetica Neue"
MONO  = "Menlo"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── helpers ──────────────────────────────────────────────────────────────────
def new_slide(prs: Presentation, bg: RGBColor = CREAM):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_rect.line.fill.background()
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = bg
    return slide


def add_text(
    slide,
    x, y, w, h,
    text: str,
    *,
    font=SERIF,
    size=14,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
    tracking=None,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, *, color=HAIRLINE, weight=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def hairline(slide, x, y, w, *, color=HAIRLINE, weight=0.75):
    add_line(slide, x, y, x + w, y, color=color, weight=weight)


def page_chrome(slide, page_num: int, total: int = 16, dark: bool = False):
    """Eyebrow + page number on every interior slide."""
    label_col = CREAM_ON_D if dark else MUTED
    # eyebrow text top-left
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
        "BIOSHIELD  ·  CLASS PRESENTATION",
        font=SANS, size=9, color=label_col, tracking=True,
    )
    # page number top-right
    add_text(
        slide, Inches(11.7), Inches(0.4), Inches(1.1), Inches(0.3),
        f"{page_num:02d} / {total:02d}",
        font=SANS, size=9, color=label_col, align=PP_ALIGN.RIGHT,
    )
    # subtle rule under header
    rule_color = HAIRLINE if not dark else RGBColor(0x45, 0x4F, 0x5E)
    hairline(slide, Inches(0.6), Inches(0.75), Inches(12.1), color=rule_color)


# ─── slide 1: title ───────────────────────────────────────────────────────────
def slide_title(prs):
    s = new_slide(prs, CREAM)
    # thin coral accent at top
    add_rect(s, Inches(0.6), Inches(0.6), Inches(0.5), Inches(0.04), fill=CORAL)
    add_text(
        s, Inches(0.6), Inches(0.78), Inches(8), Inches(0.3),
        "CLASS PRESENTATION  ·  APRIL 2026",
        font=SANS, size=10, color=MUTED,
    )
    # big serif wordmark
    add_text(
        s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.6),
        "BioShield",
        font=SERIF, size=96, color=INK, italic=True,
    )
    # tagline
    add_text(
        s, Inches(0.6), Inches(3.6), Inches(11.5), Inches(1.2),
        "Teaching a machine to spot fake\nmedical research — before it spreads.",
        font=SERIF, size=34, color=SOFT_INK, line_spacing=1.25,
    )
    # footer author block
    hairline(s, Inches(0.6), Inches(6.6), Inches(12.1))
    add_text(
        s, Inches(0.6), Inches(6.75), Inches(8), Inches(0.4),
        "Venkatarayudu Alapati  ·  AP25122040029",
        font=SANS, size=11, color=INK,
    )
    add_text(
        s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "M.Tech AI/ML  ·  NLP course project",
        font=SANS, size=10, color=MUTED,
    )
    add_text(
        s, Inches(9.2), Inches(6.75), Inches(3.6), Inches(0.4),
        "github.com/Rayudualapati-25/BioShield",
        font=MONO, size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ─── slide 2: opening hook ────────────────────────────────────────────────────
def slide_hook(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 2)
    # centered serif quote
    add_text(
        s, Inches(1.2), Inches(2.3), Inches(10.9), Inches(2.5),
        "“One of these abstracts\nwas written by a person.\nThe other was written by a\n7-billion-parameter model.”",
        font=SERIF, size=42, color=INK, italic=True,
        align=PP_ALIGN.CENTER, line_spacing=1.3,
    )
    add_text(
        s, Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.6),
        "You have thirty seconds.",
        font=SANS, size=18, color=CORAL, align=PP_ALIGN.CENTER,
    )


# ─── slide 3: the quiz ────────────────────────────────────────────────────────
def slide_quiz(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 3)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
        "The quiz.",
        font=SERIF, size=36, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.35),
        "Which abstract was written by a human?",
        font=SANS, size=13, color=MUTED,
    )

    # two cards
    card_y = Inches(2.3)
    card_h = Inches(4.4)
    card_w = Inches(5.8)
    gap_x = Inches(0.5)
    left_x = Inches(0.6)
    right_x = left_x + card_w + gap_x

    real_text = (
        "pheochromocytoma is a catecholamine-producing tumor. there are a "
        "very few reported cases of clinical pheochromocytoma. here, we "
        "report a 27-year-old woman para 1 live 1 with chief complaint of "
        "headache, confusion, nausea, and vomiting 2 days after cesarean "
        "section. she was anxious and had palpitation. on physical "
        "examination, fever, tachycardia, tachypnea, high blood pressure, "
        "and diaphoresis were observed…"
    )
    fake_text = (
        "Objective: To examine the efficacy of liraglutide and semaglutide "
        "in managing obesity. Methods: We conducted a systematic review of "
        "the PubMed database with the keywords “obesity”, “liraglutide”, "
        "“semaglutide”, and “clinical trial”. The results of these trials "
        "were combined using meta-analysis. Results: The results showed "
        "that liraglutide and semaglutide were effective in managing "
        "obesity…"
    )

    for x, letter, body in [(left_x, "A", real_text), (right_x, "B", fake_text)]:
        add_rect(s, x, card_y, card_w, card_h, fill=PAPER, line=CARD_EDGE, line_w=Pt(0.5))
        # letter badge
        add_text(
            s, x + Inches(0.3), card_y + Inches(0.2), Inches(0.8), Inches(0.8),
            letter,
            font=SERIF, size=40, color=CORAL, italic=True,
        )
        # body text
        add_text(
            s, x + Inches(0.35), card_y + Inches(1.1), card_w - Inches(0.7), card_h - Inches(1.3),
            body,
            font=SERIF, size=13, color=INK, line_spacing=1.35,
        )

    add_text(
        s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.3),
        "(We'll come back to this in five slides.)",
        font=SANS, size=10, color=MUTED, italic=True,
    )


# ─── slide 4: why this matters ────────────────────────────────────────────────
def slide_why(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 4)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
        "Why this matters.",
        font=SERIF, size=36, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
        "Three ways fake medical writing causes real damage.",
        font=SANS, size=14, color=MUTED,
    )

    panels = [
        ("Misinformation", "A plausible-looking fake abstract on vaccine safety or drug efficacy can reach millions before anyone reads the methods section."),
        ("Database rot",   "PubMed, CORD-19, and other corpora are ingested by downstream tools. Polluted sources mean polluted models, systematic reviews, and meta-analyses."),
        ("Research fraud", "Paper mills now generate synthetic manuscripts and reviews at scale. Journals need a screening layer that doesn't rely on a human's pattern-matching."),
    ]
    panel_y = Inches(3.0)
    panel_h = Inches(3.6)
    panel_w = Inches(3.85)
    gap = Inches(0.3)
    start_x = Inches(0.6)

    for i, (title, body) in enumerate(panels):
        x = start_x + (panel_w + gap) * i
        # hairline above each panel
        hairline(s, x, panel_y, panel_w, color=CORAL, weight=1.5)
        add_text(
            s, x, panel_y + Inches(0.2), panel_w, Inches(0.5),
            f"0{i+1}",
            font=SANS, size=11, color=CORAL, bold=True,
        )
        add_text(
            s, x, panel_y + Inches(0.6), panel_w, Inches(0.7),
            title,
            font=SERIF, size=22, color=INK, italic=True,
        )
        add_text(
            s, x, panel_y + Inches(1.5), panel_w, panel_h - Inches(1.5),
            body,
            font=SERIF, size=13, color=SOFT_INK, line_spacing=1.45,
        )


# ─── slide 5: approach ────────────────────────────────────────────────────────
def slide_approach(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 5)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
        "Three models in an adversarial loop.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
        "The idea: let a generator forge fakes, a detector catch them, and a coach make the fakes harder.",
        font=SANS, size=13, color=MUTED,
    )

    # three circles horizontally
    nodes = [
        ("Generator",  "BioMistral-7B\nwrites convincing fakes", CORAL),
        ("Detector",   "BiomedBERT-Large\nreal vs fake binary", INK),
        ("Agent",      "Qwen2.5-7B\nrewrites near-misses", FOREST),
    ]
    node_y = Inches(3.8)
    node_d = Inches(1.8)
    label_y = node_y + node_d + Inches(0.15)
    positions_x = [Inches(1.2), Inches(5.75), Inches(10.3)]

    for (name, body, color), x in zip(nodes, positions_x):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x, node_y, node_d, node_d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = PAPER
        circ.line.color.rgb = color
        circ.line.width = Pt(2.25)
        circ.shadow.inherit = False
        # node title centered inside circle
        add_text(
            s, x, node_y + Inches(0.65), node_d, Inches(0.5),
            name,
            font=SERIF, size=20, color=color, italic=True,
            align=PP_ALIGN.CENTER,
        )
        # body under node
        add_text(
            s, x - Inches(0.6), label_y, node_d + Inches(1.2), Inches(1.1),
            body,
            font=SANS, size=12, color=SOFT_INK,
            align=PP_ALIGN.CENTER, line_spacing=1.35,
        )

    # arrows between nodes
    arrow_y = node_y + node_d / 2
    for sx, ex in [(positions_x[0] + node_d, positions_x[1]),
                   (positions_x[1] + node_d, positions_x[2])]:
        conn = s.shapes.add_connector(1, sx + Inches(0.15), arrow_y, ex - Inches(0.15), arrow_y)
        conn.line.color.rgb = MUTED
        conn.line.width = Pt(1.25)

    # pull-quote line under diagram
    add_text(
        s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
        "Each round, the detector is retrained on harder fakes. That's the loop.",
        font=SERIF, size=15, color=SOFT_INK, italic=True, align=PP_ALIGN.CENTER,
    )


# ─── slide 6: the cast ────────────────────────────────────────────────────────
def slide_cast(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 6)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
        "The cast.",
        font=SERIF, size=36, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
        "Five models. Each picked for a specific job.",
        font=SANS, size=13, color=MUTED,
    )

    rows = [
        ("BioMistral-7B",   "The forger",      "Medical-tuned Mistral. Zero-shot generation of fake abstracts."),
        ("BiomedBERT-Large", "The skeptic",    "340M-param encoder pretrained on biomedical text. Binary classifier."),
        ("Qwen2.5-7B",      "The coach",       "Rewrites the fakes that almost fooled the detector. Picks sharper wording."),
        ("SeqGAN",          "The scrappy challenger", "Small GAN. Used in Condition B to produce stylistically-different fakes."),
        ("Llama-3.2-3B",    "The control",     "Held back for cross-family transfer-attack evaluation."),
    ]
    row_y = Inches(2.7)
    row_h = Inches(0.78)
    for i, (name, role, desc) in enumerate(rows):
        y = row_y + row_h * i
        # hairline between rows
        hairline(s, Inches(0.6), y, Inches(12.1))
        add_text(
            s, Inches(0.7), y + Inches(0.15), Inches(3.4), Inches(0.5),
            name,
            font=SERIF, size=18, color=INK, italic=True,
        )
        add_text(
            s, Inches(4.2), y + Inches(0.2), Inches(2.5), Inches(0.4),
            role,
            font=SANS, size=11, color=CORAL, bold=True,
        )
        add_text(
            s, Inches(6.8), y + Inches(0.22), Inches(6.0), Inches(0.5),
            desc,
            font=SERIF, size=13, color=SOFT_INK,
        )
    # bottom rule
    hairline(s, Inches(0.6), row_y + row_h * len(rows), Inches(12.1))


# ─── slide 7: data — hero number ──────────────────────────────────────────────
def slide_data(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 7)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
        "The dataset.",
        font=SERIF, size=36, color=INK, italic=True,
    )
    # hero numeral
    add_text(
        s, Inches(0.6), Inches(2.3), Inches(6.5), Inches(3.2),
        "1,000",
        font=SERIF, size=180, color=INK, italic=True, line_spacing=1.0,
    )
    add_text(
        s, Inches(0.6), Inches(5.6), Inches(6.5), Inches(0.45),
        "biomedical abstracts",
        font=SERIF, size=22, color=SOFT_INK, italic=True,
    )
    # right-column breakdown
    rx = Inches(8.0)
    add_text(
        s, rx, Inches(2.5), Inches(4.8), Inches(0.3),
        "HOW IT BREAKS DOWN",
        font=SANS, size=10, color=CORAL, bold=True,
    )
    hairline(s, rx, Inches(2.85), Inches(4.8), color=CORAL)
    breakdown = [
        ("500", "real PubMed abstracts, streamed from ccdv/pubmed-summarization"),
        ("500", "fake abstracts, generated zero-shot by BioMistral-7B"),
        ("80 / 10 / 10", "train / val / test split, stratified by label"),
        ("n = 100", "held-out test set, same distribution as training"),
    ]
    for i, (num, desc) in enumerate(breakdown):
        y = Inches(3.1) + Inches(0.85) * i
        add_text(s, rx, y, Inches(4.8), Inches(0.4), num,
                 font=SERIF, size=20, color=INK, italic=True)
        add_text(s, rx, y + Inches(0.42), Inches(4.8), Inches(0.4), desc,
                 font=SANS, size=11, color=SOFT_INK)


# ─── slide 8: four experiments ────────────────────────────────────────────────
def slide_experiments(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 8)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
        "Four experiments, one ablation.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
        "Changing one variable at a time tells us which piece actually earns its place.",
        font=SANS, size=13, color=MUTED,
    )

    cards = [
        ("A", "static_baseline", "Train once. Never retrain.",
         "The reference point — how good is a detector without any adversarial pressure?"),
        ("B", "seqgan_only",     "SeqGAN + retrain each round.",
         "Pure adversarial dynamic. No coach, no fancy rewrites — just new fakes each round."),
        ("C", "agent_only",      "Qwen rewrites, no retrain.",
         "Does a smarter attacker alone hurt a static detector? Rewrite the hard fakes and measure."),
        ("D", "full_pipeline",   "SeqGAN + Qwen + retrain.",
         "Everything on. The setup that should, in principle, produce the hardest detector."),
    ]
    # 2×2 grid
    card_w = Inches(5.9)
    card_h = Inches(2.3)
    gap = Inches(0.3)
    x0 = Inches(0.6)
    y0 = Inches(2.65)
    for idx, (letter, name, short, desc) in enumerate(cards):
        col, row = idx % 2, idx // 2
        x = x0 + (card_w + gap) * col
        y = y0 + (card_h + gap) * row
        add_rect(s, x, y, card_w, card_h, fill=PAPER, line=CARD_EDGE, line_w=Pt(0.5))
        # big letter
        add_text(
            s, x + Inches(0.25), y + Inches(0.15), Inches(1.0), Inches(1.2),
            letter,
            font=SERIF, size=56, color=CORAL, italic=True,
        )
        add_text(
            s, x + Inches(1.3), y + Inches(0.3), card_w - Inches(1.4), Inches(0.5),
            name,
            font=MONO, size=13, color=MUTED,
        )
        add_text(
            s, x + Inches(1.3), y + Inches(0.7), card_w - Inches(1.4), Inches(0.55),
            short,
            font=SERIF, size=19, color=INK, italic=True,
        )
        add_text(
            s, x + Inches(1.3), y + Inches(1.35), card_w - Inches(1.4), card_h - Inches(1.4),
            desc,
            font=SERIF, size=12, color=SOFT_INK, line_spacing=1.35,
        )


# ─── slide 9: result 1 — AUC hero ─────────────────────────────────────────────
def slide_result_auc(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 9)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "RESULT 01",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "The baseline detector is already strong.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    # hero numeral
    add_text(
        s, Inches(0.6), Inches(2.6), Inches(9), Inches(3.2),
        "0.989",
        font=SERIF, size=220, color=INK, italic=True, line_spacing=1.0,
    )
    add_text(
        s, Inches(0.6), Inches(5.95), Inches(9), Inches(0.45),
        "Test-set AUC  ·  Condition A  ·  n = 100",
        font=SANS, size=14, color=SOFT_INK,
    )
    # right: supporting numbers
    rx = Inches(9.4)
    add_text(s, rx, Inches(2.7), Inches(3.4), Inches(0.3), "ALSO WORTH NOTING",
             font=SANS, size=10, color=CORAL, bold=True)
    hairline(s, rx, Inches(3.05), Inches(3.4), color=CORAL)
    mini = [("F1", "0.938"), ("Accuracy", "0.94"), ("Train samples", "800"), ("Epochs", "2")]
    for i, (k, v) in enumerate(mini):
        y = Inches(3.25) + Inches(0.75) * i
        add_text(s, rx, y, Inches(1.8), Inches(0.4), k, font=SANS, size=11, color=MUTED)
        add_text(s, rx + Inches(1.8), y - Inches(0.05), Inches(1.6), Inches(0.5), v,
                 font=SERIF, size=22, color=INK, italic=True, align=PP_ALIGN.RIGHT)


# ─── slide 10: result 2 — evasion curve ───────────────────────────────────────
def slide_result_evasion(prs, evasion_plot: Path | None):
    s = new_slide(prs, CREAM)
    page_chrome(s, 10)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "RESULT 02",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "The detector is learning.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.25), Inches(12), Inches(0.45),
        "Condition B  ·  evasion rate on the SeqGAN fake pool, round over round.",
        font=SANS, size=13, color=MUTED,
    )

    # three hero numbers with arrows
    numbers = [("0.74", "ROUND 1"), ("0.64", "ROUND 2"), ("0.54", "ROUND 3")]
    nx0 = Inches(0.7)
    nbox_w = Inches(3.5)
    for i, (num, label) in enumerate(numbers):
        x = nx0 + (nbox_w + Inches(0.3)) * i
        add_text(
            s, x, Inches(3.1), nbox_w, Inches(1.6),
            num,
            font=SERIF, size=86, color=(CORAL if i == 0 else (MUTED if i == 1 else FOREST)),
            italic=True, align=PP_ALIGN.CENTER,
        )
        add_text(
            s, x, Inches(4.75), nbox_w, Inches(0.4),
            label,
            font=SANS, size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=True,
        )
        # arrow
        if i < 2:
            ax = x + nbox_w + Inches(0.02)
            ay = Inches(3.9)
            conn = s.shapes.add_connector(1, ax, ay, ax + Inches(0.25), ay)
            conn.line.color.rgb = MUTED
            conn.line.width = Pt(1.25)

    # interpretation block
    add_rect(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.3),
             fill=CORAL_SOFT)
    add_text(
        s, Inches(0.85), Inches(5.75), Inches(11.5), Inches(1.2),
        "That's a 27% relative reduction in three rounds. The detector is seeing fakes\n"
        "it previously missed, retraining on them, and getting harder to fool. This is\n"
        "the adversarial-training story in one graph.",
        font=SERIF, size=15, color=INK, italic=True, line_spacing=1.35,
    )


# ─── slide 11: result 3 — what didn't work ────────────────────────────────────
def slide_result_honest(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 11)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "RESULT 03",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "And here's what didn't work.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.5),
        "Conditions C and D produced byte-identical final metrics.",
        font=SERIF, size=22, color=SOFT_INK, italic=True, line_spacing=1.35,
    )
    # two side-by-side plain rows
    y = Inches(3.4)
    rows = [
        ("Condition C", "agent_only", "AUC 0.982  ·  F1 0.947  ·  Evasion 0.00"),
        ("Condition D", "full_pipeline", "AUC 0.982  ·  F1 0.947  ·  Evasion 0.00"),
    ]
    for i, (label, code, vals) in enumerate(rows):
        ry = y + Inches(0.7) * i
        hairline(s, Inches(0.6), ry, Inches(12.1))
        add_text(s, Inches(0.7), ry + Inches(0.1), Inches(2.5), Inches(0.45),
                 label, font=SERIF, size=16, color=INK, italic=True)
        add_text(s, Inches(3.2), ry + Inches(0.15), Inches(2.8), Inches(0.4),
                 code, font=MONO, size=11, color=MUTED)
        add_text(s, Inches(6.5), ry + Inches(0.12), Inches(6.2), Inches(0.45),
                 vals, font=MONO, size=12, color=INK, align=PP_ALIGN.RIGHT)
    hairline(s, Inches(0.6), y + Inches(0.7) * len(rows), Inches(12.1))

    # explanation
    add_text(
        s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.8),
        "Both conditions picked Condition A's baseline as their best-validation checkpoint.\n"
        "The cause, once you look: BioMistral's zero-shot fakes are too easy — the detector catches nearly\n"
        "all of them on the first pass. With nothing slipping through, there's no gradient for the agent\n"
        "or the retrain loop to improve against.",
        font=SERIF, size=14, color=SOFT_INK, line_spacing=1.5,
    )
    # lesson line
    add_text(
        s, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4),
        "Lesson: adversarial setups need adversarial-quality fakes.",
        font=SERIF, size=15, color=CORAL, italic=True, bold=True,
    )


# ─── slide 12: result 4 — RAID scope ──────────────────────────────────────────
def slide_result_scope(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 12)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "RESULT 04",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "What we can't claim.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.6), Inches(7.5), Inches(3.5),
        "96%",
        font=SERIF, size=220, color=CORAL, italic=True, line_spacing=1.0,
    )
    add_text(
        s, Inches(0.6), Inches(5.95), Inches(8), Inches(0.5),
        "of cross-domain fakes slipped past our detector on the RAID benchmark.",
        font=SERIF, size=15, color=SOFT_INK, italic=True, line_spacing=1.35,
    )
    # right column — framing
    rx = Inches(8.4)
    add_text(s, rx, Inches(2.7), Inches(4.3), Inches(0.3),
             "READ THIS HONESTLY", font=SANS, size=10, color=CORAL, bold=True)
    hairline(s, rx, Inches(3.05), Inches(4.3), color=CORAL)
    add_text(
        s, rx, Inches(3.2), Inches(4.3), Inches(3.5),
        "RAID contains general-web AI text from ChatGPT, GPT-4, Claude, Llama, and others — "
        "models our detector never trained against.\n\n"
        "BioShield is a biomedical specialist, not a universal AI-text detector. That's a "
        "positioning statement, not a weakness.\n\n"
        "Cross-domain generalization is its own research problem.",
        font=SERIF, size=13, color=SOFT_INK, line_spacing=1.5,
    )


# ─── slide 13: under the hood — bugs ──────────────────────────────────────────
def slide_bugs(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 13)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "ENGINEERING",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "Two bugs we found in our own code.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
        "Both caught during an earlier smoke run. Both verified fixed on the 500-scale run.",
        font=SANS, size=13, color=MUTED,
    )

    bugs = [
        ("01",
         "The seeding bug",
         "set_seed(42) was called once at process start — so every adversarial round produced byte-identical fakes and byte-identical rewrites.",
         "Before:  rounds 1, 2, 3 had literally the same output.\nAfter:   set_seed(base + round_idx) inside the loop — rounds now vary as intended."),
        ("02",
         "The hard-fraction bug",
         "The config had hard_fraction=0.25, but the pipeline ignored it and sent all 250 fakes to Qwen each round.",
         "Before:  Qwen rewrote everything.\nAfter:   sort by detector p(real) → keep top 25% hardest → Qwen rewrites those only.\nVerified live: round_1/fakes_hard_input.csv, fakes_hard_rewritten.csv on disk."),
    ]
    card_y = Inches(3.0)
    card_h = Inches(3.7)
    card_w = Inches(5.95)
    gap = Inches(0.25)
    for i, (num, title, problem, fix) in enumerate(bugs):
        x = Inches(0.6) + (card_w + gap) * i
        add_rect(s, x, card_y, card_w, card_h, fill=PAPER, line=CARD_EDGE, line_w=Pt(0.5))
        add_text(s, x + Inches(0.35), card_y + Inches(0.25), Inches(1), Inches(0.4),
                 num, font=SANS, size=10, color=CORAL, bold=True)
        add_text(s, x + Inches(0.35), card_y + Inches(0.55), card_w - Inches(0.7), Inches(0.6),
                 title, font=SERIF, size=22, color=INK, italic=True)
        add_text(s, x + Inches(0.35), card_y + Inches(1.3), card_w - Inches(0.7), Inches(1.1),
                 problem, font=SERIF, size=13, color=SOFT_INK, line_spacing=1.45)
        add_text(s, x + Inches(0.35), card_y + Inches(2.5), card_w - Inches(0.7), Inches(1.3),
                 fix, font=MONO, size=10, color=INK, line_spacing=1.45)


# ─── slide 14: runtime ────────────────────────────────────────────────────────
def slide_runtime(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 14)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "INFRASTRUCTURE",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "Six hours and twenty-three minutes.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
        "On a single M3 Max — no cluster, no cloud. Here's where the time went.",
        font=SANS, size=13, color=MUTED,
    )

    # horizontal stacked bar, proportional
    bar_x = Inches(0.6)
    bar_y = Inches(3.3)
    bar_w = Inches(12.1)
    bar_h = Inches(0.55)
    segments = [
        ("BioMistral generation", 0.70, CORAL),
        ("Detector retrain",       0.20, INK),
        ("Qwen rewrites",          0.08, FOREST),
        ("Everything else",        0.02, MUTED),
    ]
    cursor = bar_x
    for label, frac, color in segments:
        seg_w = Emu(int(bar_w * frac))
        add_rect(s, cursor, bar_y, seg_w, bar_h, fill=color)
        cursor += seg_w

    # legend underneath
    legend_y = Inches(4.1)
    lx = Inches(0.6)
    for i, (label, frac, color) in enumerate(segments):
        yi = legend_y + Inches(0.5) * i
        swatch = add_rect(s, lx, yi + Inches(0.08), Inches(0.18), Inches(0.18), fill=color)
        add_text(s, lx + Inches(0.35), yi, Inches(6), Inches(0.35),
                 label, font=SANS, size=13, color=INK)
        add_text(s, lx + Inches(6.4), yi, Inches(1.5), Inches(0.35),
                 f"{int(frac * 100)}%", font=SERIF, size=16, color=INK, italic=True, align=PP_ALIGN.LEFT)

    # right-column facts
    rx = Inches(8.5)
    add_text(s, rx, Inches(4.1), Inches(4.2), Inches(0.3),
             "DEVICE NOTES", font=SANS, size=10, color=CORAL, bold=True)
    hairline(s, rx, Inches(4.45), Inches(4.2), color=CORAL)
    facts = [
        ("M3 Max 64 GB", "single-machine, lid closed via caffeinate"),
        ("bf16 on MPS",  "no fp16 underflow, no cluster needed"),
        ("3,984 samples produced", "across real, BioMistral, SeqGAN, RAID"),
    ]
    for i, (k, v) in enumerate(facts):
        yi = Inches(4.65) + Inches(0.75) * i
        add_text(s, rx, yi, Inches(4.2), Inches(0.4), k,
                 font=SERIF, size=16, color=INK, italic=True)
        add_text(s, rx, yi + Inches(0.4), Inches(4.2), Inches(0.4), v,
                 font=SANS, size=11, color=SOFT_INK)


# ─── slide 15: paper roadmap ──────────────────────────────────────────────────
def slide_paper(prs):
    s = new_slide(prs, CREAM)
    page_chrome(s, 15)
    add_text(
        s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
        "WHAT'S NEXT",
        font=SANS, size=11, color=CORAL, bold=True,
    )
    add_text(
        s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.7),
        "From class project to paper.",
        font=SERIF, size=34, color=INK, italic=True,
    )
    add_text(
        s, Inches(0.6), Inches(2.35), Inches(12), Inches(0.4),
        "The numbers we have are project-worthy. Publishable needs tighter confidence intervals and harder fakes.",
        font=SANS, size=13, color=MUTED,
    )

    # three-column comparison table
    cols = [
        ("NOW",        "Class presentation",   "1,000 samples",  "n = 100",     "6 h 23 m",  "±5.0% CI", MUTED),
        ("TIER 1",     "Pilot — peer-review defensible", "4,000 samples",  "n = 200",     "~24 h",     "±3.5% CI", INK),
        ("TIER 2",     "Full — reviewer-proof", "20,000 samples", "n = 1,000",   "~80 h",     "±1.5% CI", CORAL),
    ]
    cw = Inches(3.95)
    cx0 = Inches(0.6)
    gap = Inches(0.25)
    cy = Inches(3.1)
    ch = Inches(3.2)
    for i, (label, name, data, test, wall, ci, color) in enumerate(cols):
        x = cx0 + (cw + gap) * i
        # card
        add_rect(s, x, cy, cw, ch, fill=PAPER, line=CARD_EDGE, line_w=Pt(0.5))
        # colored strip at top
        add_rect(s, x, cy, cw, Inches(0.08), fill=color)
        add_text(s, x + Inches(0.3), cy + Inches(0.25), cw - Inches(0.6), Inches(0.35),
                 label, font=SANS, size=11, color=color, bold=True)
        add_text(s, x + Inches(0.3), cy + Inches(0.6), cw - Inches(0.6), Inches(0.5),
                 name, font=SERIF, size=17, color=INK, italic=True)
        facts = [("DATA", data), ("TEST SET", test), ("WALL-CLOCK", wall), ("95% CI", ci)]
        for j, (k, v) in enumerate(facts):
            yf = cy + Inches(1.3) + Inches(0.44) * j
            add_text(s, x + Inches(0.3), yf, Inches(1.6), Inches(0.3),
                     k, font=SANS, size=9, color=MUTED, bold=True)
            add_text(s, x + Inches(1.8), yf - Inches(0.03), cw - Inches(2.0), Inches(0.35),
                     v, font=SERIF, size=14, color=INK, italic=True, align=PP_ALIGN.RIGHT)

    # footer bullet list
    add_text(
        s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.4),
        "Pre-work  ·  batch BioMistral gen  ·  LoRA-hardened generator  ·  replace M4 dataset  ·  populate transfer_test  ·  human eval",
        font=SANS, size=11, color=SOFT_INK, italic=True,
    )


# ─── slide 16: takeaways + questions ──────────────────────────────────────────
def slide_takeaways(prs):
    s = new_slide(prs, DEEP)  # dramatic dark background for closer
    # darker rule colors
    add_rect(s, Inches(0.6), Inches(0.6), Inches(0.5), Inches(0.04), fill=CORAL)
    add_text(
        s, Inches(0.6), Inches(0.78), Inches(8), Inches(0.3),
        "CLASS PRESENTATION  ·  APRIL 2026",
        font=SANS, size=10, color=CREAM_ON_D,
    )
    add_text(
        s, Inches(11.7), Inches(0.78), Inches(1.1), Inches(0.3),
        "16 / 16",
        font=SANS, size=10, color=CREAM_ON_D, align=PP_ALIGN.RIGHT,
    )

    add_text(
        s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.9),
        "Three things to take with you.",
        font=SERIF, size=40, color=CREAM_ON_D, italic=True,
    )
    # three numbered statements
    items = [
        ("01",
         "Adversarial retraining works — when fakes are actually hard.",
         "Condition B: 27% relative evasion reduction in three rounds."),
        ("02",
         "Zero-shot fakes weren't hard enough to train against.",
         "C and D converged on A's checkpoint. Harder generators = future work."),
        ("03",
         "Specialist > universal for this problem.",
         "RAID 0.96 evasion means one detector can't guard all domains. Own that."),
    ]
    ty = Inches(2.9)
    for i, (num, headline, tail) in enumerate(items):
        y = ty + Inches(1.1) * i
        add_text(s, Inches(0.6), y, Inches(0.9), Inches(0.55),
                 num, font=SERIF, size=30, color=CORAL, italic=True)
        add_text(s, Inches(1.5), y + Inches(0.05), Inches(11), Inches(0.5),
                 headline, font=SERIF, size=19, color=CREAM_ON_D, italic=True)
        add_text(s, Inches(1.5), y + Inches(0.55), Inches(11), Inches(0.45),
                 tail, font=SANS, size=12, color=RGBColor(0xB0, 0xB0, 0xB0))

    # closing
    hairline(s, Inches(0.6), Inches(6.5), Inches(12.1),
             color=RGBColor(0x45, 0x4F, 0x5E))
    add_text(
        s, Inches(0.6), Inches(6.7), Inches(8), Inches(0.5),
        "Questions?",
        font=SERIF, size=28, color=CORAL, italic=True,
    )
    add_text(
        s, Inches(8.0), Inches(6.85), Inches(4.7), Inches(0.3),
        "Venkatarayudu Alapati  ·  AP25122040029",
        font=SANS, size=10, color=CREAM_ON_D, align=PP_ALIGN.RIGHT,
    )
    add_text(
        s, Inches(8.0), Inches(7.15), Inches(4.7), Inches(0.3),
        "github.com/Rayudualapati-25/BioShield",
        font=MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ─── orchestrator ─────────────────────────────────────────────────────────────
def build(output_path: Path, evasion_plot: Path | None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_hook(prs)
    slide_quiz(prs)
    slide_why(prs)
    slide_approach(prs)
    slide_cast(prs)
    slide_data(prs)
    slide_experiments(prs)
    slide_result_auc(prs)
    slide_result_evasion(prs, evasion_plot)
    slide_result_honest(prs)
    slide_result_scope(prs)
    slide_bugs(prs)
    slide_runtime(prs)
    slide_paper(prs)
    slide_takeaways(prs)

    prs.save(str(output_path))
    print(f"Wrote {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    project = Path(__file__).resolve().parent.parent
    out = project / "experiments" / "BioShield_class_presentation_v2.pptx"
    plot = project / "experiments" / "plots" / "evasion_rate_vs_round.png"
    build(out, plot if plot.exists() else None)
